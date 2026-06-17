"""
Generates CloudDrive_Presentation.pptx — a polished slide deck presenting the
whole project, reusing the diagrams and real screenshots from docs/img/.

Run:  python docs/generate_slides.py
"""
import os
import struct

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(__file__)
IMG = os.path.join(HERE, "img")
OUT = os.path.join(HERE, "CloudDrive_Presentation.pptx")

# ---- palette ----
BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1D, 0x4E, 0xD8)
NAVY = RGBColor(0x10, 0x23, 0x55)
INK = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY = RGBColor(0x38, 0xBD, 0xF8)
GREEN = RGBColor(0x05, 0x96, 0x69)

EMU = 914400
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_slide_no = 0


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(24)
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def _rect(slide, l, t, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _text(slide, l, t, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb


def fit_image(slide, path, box_l, box_t, box_w, box_h):
    iw, ih = png_size(path)
    scale = min(box_w / iw, box_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    l = int(box_l + (box_w - w) / 2)
    t = int(box_t + (box_h - h) / 2)
    slide.shapes.add_picture(path, l, t, width=w, height=h)
    return l, t, w, h


def content_slide(title, eyebrow=None):
    """White slide with a title bar; returns slide + content top."""
    global _slide_no
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    _slide_no += 1
    # accent bar
    _rect(s, 0, 0, Inches(0.22), SH, BLUE)
    if eyebrow:
        _text(s, Inches(0.55), Inches(0.34), Inches(11), Inches(0.4),
              eyebrow.upper(), 12, BLUE, bold=True)
        ttop = Inches(0.66)
    else:
        ttop = Inches(0.45)
    _text(s, Inches(0.5), ttop, Inches(12.3), Inches(0.9), title, 30, INK, bold=True)
    _rect(s, Inches(0.55), Inches(1.5), Inches(1.6), Inches(0.05), SKY)
    # footer
    _text(s, Inches(0.5), Inches(7.02), Inches(8), Inches(0.4),
          "CloudDrive — Scalable File Storage on AWS", 9, GREY)
    _text(s, Inches(12.2), Inches(7.02), Inches(0.9), Inches(0.4),
          str(_slide_no), 9, GREY, align=PP_ALIGN.RIGHT)
    return s


def bullets(slide, items, l, t, w, h, size=17, gap=10, color=INK):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        lvl = 0
        txt = it
        if isinstance(it, tuple):
            txt, lvl = it
        r = p.add_run()
        r.text = ("•   " if lvl == 0 else "–   ") + txt
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = color if lvl == 0 else GREY
        r.font.name = "Calibri"
        p.level = lvl
    return tb


# ======================================================================= #
#  SLIDE 1 — TITLE
# ======================================================================= #
def title_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    _rect(s, 0, 0, SW, SH, NAVY)
    # decorative band
    _rect(s, 0, Inches(4.7), SW, Inches(0.08), BLUE)
    # logo tile
    tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.5),
                              Inches(1.5), Inches(1.5))
    tile.fill.solid(); tile.fill.fore_color.rgb = WHITE; tile.line.fill.background()
    tile.shadow.inherit = False
    _text(s, Inches(0.9), Inches(1.62), Inches(1.5), Inches(1.3), "☁", 60, BLUE,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(2.7), Inches(1.55), Inches(10), Inches(1.4), "CloudDrive", 60, WHITE, bold=True)
    _text(s, Inches(2.78), Inches(2.95), Inches(10), Inches(0.7),
          "Scalable File Storage on AWS — a Mini Google Drive", 22, SKY)
    _text(s, Inches(0.95), Inches(5.0), Inches(11), Inches(0.5),
          "Final-Year Cloud Computing Project", 16, WHITE)
    _text(s, Inches(0.95), Inches(5.7), Inches(11.5), Inches(1.2),
          "Team:  Shehriyar Ali Rustam  ·  Touseef Abbas  ·  Abdul Ahad", 15, RGBColor(0xC7, 0xD2, 0xFE))
    _text(s, Inches(0.95), Inches(6.5), Inches(11), Inches(0.4),
          "github.com/Shehriyar-Ali-Rustam/clouddrive", 12, RGBColor(0x93, 0xA5, 0xCF))


# ======================================================================= #
#  BUILD DECK
# ======================================================================= #
def img(p):
    return os.path.join(IMG, p)


def build():
    title_slide()

    # 2 — Problem & Solution
    s = content_slide("The Problem & Our Solution", "Motivation")
    bullets(s, [
        "Storing and sharing files across devices and people is a universal need.",
        "Traditional apps stream every file through the server — it becomes a bottleneck and a single point of failure.",
        ("So it cannot scale to many users or large files.", 1),
        "CloudDrive solves this with cloud object storage (Amazon S3) and pre-signed URLs.",
        ("The browser sends file bytes directly to S3 — the server stays light and scalable.", 1),
    ], Inches(0.6), Inches(1.9), Inches(12), Inches(4.8), size=18, gap=12)

    # 3 — What is CloudDrive (features + dashboard)
    s = content_slide("What is CloudDrive?", "Overview")
    bullets(s, [
        "A secure, cloud-based 'Mini Google Drive'.",
        "Sign up, upload, organise, download files.",
        "Share with users or via public expiring links.",
        "Per-user storage quotas & isolation.",
        "Runs on real AWS S3; deployable via Docker.",
    ], Inches(0.6), Inches(1.95), Inches(5.4), Inches(4.6), size=18, gap=12)
    fit_image(s, img("02_dashboard.png"), Inches(6.2), Inches(1.85), Inches(6.7), Inches(4.7))

    # 4 — Architecture
    s = content_slide("System Architecture", "How it fits together")
    fit_image(s, img("diag_arch.png"), Inches(0.6), Inches(1.75), Inches(8.1), Inches(5.0))
    bullets(s, [
        "Browser → CDN/Load Balancer → stateless API.",
        "API on ECS Fargate (autoscaling).",
        "Metadata in PostgreSQL (RDS).",
        "File bytes in Amazon S3.",
        ("Browser talks to S3 directly via pre-signed URLs.", 1),
    ], Inches(8.9), Inches(2.0), Inches(4.1), Inches(4.5), size=15, gap=10)

    # 5 — Core idea: pre-signed URLs
    s = content_slide("The Core Idea — Pre-signed URLs", "Why it scales")
    fit_image(s, img("diag_upload.png"), Inches(0.6), Inches(1.9), Inches(12.1), Inches(2.4))
    bullets(s, [
        "The API never streams file bytes — it issues a short-lived signed 'permission slip'.",
        "The browser uploads/downloads straight to Amazon S3.",
        "Result: a stateless API you can run in many copies and autoscale — the heart of cloud design.",
    ], Inches(0.7), Inches(4.6), Inches(12), Inches(2.2), size=17, gap=12)

    # 6 — Live on AWS S3
    s = content_slide("Running on Real Amazon S3", "Cloud storage — live")
    fit_image(s, img("s3_bucket.png"), Inches(0.6), Inches(1.8), Inches(8.2), Inches(5.0))
    bullets(s, [
        "Files stored as objects in S3.",
        "Bucket versioning enabled.",
        "AES-256 encryption at rest.",
        "All public access blocked.",
        ("Verified live via the AWS CLI.", 1),
    ], Inches(9.0), Inches(2.0), Inches(4.0), Inches(4.5), size=15, gap=11)

    # 7 — Secure sharing
    s = content_slide("Secure Sharing & Collaboration", "Feature")
    fit_image(s, img("03_share_modal.png"), Inches(0.6), Inches(1.85), Inches(6.3), Inches(4.7))
    bullets(s, [
        "Share a file with a registered user by email.",
        "Generate a public link with optional expiry.",
        "A 'Shared with me' view for received files.",
        "Downloads use short-lived pre-signed URLs scoped to one file.",
        "Public links redirect straight to storage.",
    ], Inches(7.2), Inches(2.0), Inches(5.7), Inches(4.6), size=16, gap=12)

    # 8 — Tech stack
    s = content_slide("Technology Stack", "What we used")
    cols = [
        ("Backend", ["Python 3.12", "FastAPI", "Uvicorn", "SQLAlchemy", "Pydantic"]),
        ("Storage & Data", ["Amazon S3", "boto3", "PostgreSQL", "SQLite"]),
        ("Security", ["JWT (python-jose)", "bcrypt", "AWS IAM"]),
        ("Cloud / DevOps", ["Terraform", "Docker", "ECS · ALB · VPC · ECR", "GitHub Actions", "Render · Cloudflare"]),
    ]
    x = Inches(0.6)
    for name, items in cols:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(2.95), Inches(4.5))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = BLUE
        card.shadow.inherit = False
        _text(s, x, Inches(2.05), Inches(2.95), Inches(0.5), name, 16, BLUE, bold=True, align=PP_ALIGN.CENTER)
        bullets(s, items, x + Inches(0.2), Inches(2.7), Inches(2.6), Inches(3.5), size=13, gap=8)
        x += Inches(3.07)

    # 9 — Security
    s = content_slide("Security Model", "Defence in depth")
    bullets(s, [
        "Passwords hashed with bcrypt — never stored in plain text.",
        "Stateless authentication with signed, expiring JWTs.",
        "Pre-signed URLs: time-limited (15 min) and scoped to one object.",
        "Private S3 bucket, encrypted at rest (AES-256).",
        "IAM least-privilege — the app touches only its own bucket.",
        "Per-user isolation — every query filtered by owner.",
    ], Inches(0.6), Inches(1.95), Inches(12), Inches(4.7), size=18, gap=12)

    # 10 — Docker
    s = content_slide("Containerized with Docker", "Portable & deployable")
    fit_image(s, img("docker.png"), Inches(0.6), Inches(1.8), Inches(8.2), Inches(5.0))
    bullets(s, [
        "One Dockerfile packages the whole app.",
        "Built image: clouddrive-api (395 MB).",
        "Runs identically locally and on AWS ECS.",
        ("Image stored in Amazon ECR.", 1),
        ("Verified running live on S3.", 1),
    ], Inches(9.0), Inches(2.0), Inches(4.0), Inches(4.5), size=15, gap=11)

    # 11 — CI/CD
    s = content_slide("CI/CD Pipeline", "Automated delivery")
    fit_image(s, img("gh_actions.png"), Inches(0.6), Inches(1.8), Inches(8.4), Inches(5.0))
    bullets(s, [
        "GitHub Actions on every push.",
        "CI: runs all 19 tests.",
        "Deploy: build image → ECR → ECS.",
        ("Green across all runs.", 1),
    ], Inches(9.2), Inches(2.1), Inches(3.8), Inches(4.3), size=15, gap=11)

    # 12 — Deployment
    s = content_slide("Deployment Options", "Flexible, all on S3")
    fit_image(s, img("diag_deploy.png"), Inches(1.4), Inches(2.0), Inches(10.5), Inches(2.6))
    bullets(s, [
        "Local (Docker / dev) · Render (permanent URL + managed Postgres) · AWS ECS (autoscaling, Terraform).",
        "All three share the same Amazon S3 storage.",
        "A Cloudflare tunnel exposes a public HTTPS URL for live demos.",
    ], Inches(0.7), Inches(4.9), Inches(12), Inches(2.0), size=16, gap=11)

    # 13 — Testing
    s = content_slide("Testing & Quality", "Confidence")
    bullets(s, [
        "19 automated tests (pytest) — all passing.",
        ("Authentication — signup, login, wrong password, duplicate email (6).", 1),
        ("Files — upload lifecycle, quota, delete, user isolation (7).", 1),
        ("Sharing & security — user share, public link, tampered/expired URL (6).", 1),
        "Run automatically in CI on every push.",
        "Terraform validated; Docker image builds & runs.",
    ], Inches(0.6), Inches(1.95), Inches(12), Inches(4.7), size=18, gap=11)

    # 14 — Team
    s = content_slide("The Team — Who Built What", "Three modules")
    rows = [
        ("Module", "Member", "Cloud / DevOps slice"),
        ("Authentication & User Accounts", "Touseef Abbas", "Security · IAM · CI"),
        ("File Storage & Upload/Download", "Shehriyar Ali Rustam", "Amazon S3 · Terraform"),
        ("Sharing & Collaboration", "Abdul Ahad", "Docker · Deployment · Networking"),
    ]
    tbl = s.shapes.add_table(4, 3, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.4)).table
    tbl.columns[0].width = Inches(4.3); tbl.columns[1].width = Inches(3.2); tbl.columns[2].width = Inches(4.4)
    for c in range(3):
        cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = rows[0][c]
        r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = WHITE
    for ri in range(1, 4):
        for c in range(3):
            cell = tbl.cell(ri, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = rows[ri][c]
            r.font.size = Pt(14); r.font.color.rgb = INK
            if c == 1:
                r.font.bold = True; r.font.color.rgb = DARK
    _text(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.6),
          "One collective report + three individual module reports.", 14, GREY, italic=True)

    # 15 — Live demo
    s = content_slide("Live Demo", "See it running")
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(2.4), Inches(8.7), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = BLUE; box.shadow.inherit = False
    _text(s, Inches(2.3), Inches(2.55), Inches(8.7), Inches(1.0),
          "bash demo.sh  →  public HTTPS URL", 24, DARK, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, [
        "Open the link on a laptop and a phone — same app, real cloud.",
        "Upload a file → see it appear in the AWS S3 console.",
        "Create a public link → download from another device.",
        "Show the file uploading directly to S3 (DevTools) — the pre-signed URL in action.",
    ], Inches(1.8), Inches(4.2), Inches(9.7), Inches(2.4), size=17, gap=12)

    # 16 — Conclusion
    s = content_slide("Conclusion & Future Work", "Wrap-up")
    bullets(s, [
        "A complete, working, tested cloud application on real AWS S3.",
        "Demonstrates object storage, pre-signed URLs, managed databases, stateless scalable services, IaC and CI/CD.",
        "Fully reproducible from version control.",
        "Future: multipart upload for huge files; file previews & search; team workspaces.",
    ], Inches(0.6), Inches(1.95), Inches(12), Inches(4.6), size=18, gap=13)

    # 17 — Thank you
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    _rect(s, 0, Inches(3.6), SW, Inches(0.07), BLUE)
    _text(s, 0, Inches(2.4), SW, Inches(1.2), "Thank You", 54, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _text(s, 0, Inches(3.8), SW, Inches(0.7), "Questions?", 24, SKY, align=PP_ALIGN.CENTER)
    _text(s, 0, Inches(4.7), SW, Inches(0.5),
          "github.com/Shehriyar-Ali-Rustam/clouddrive", 14, RGBColor(0xC7, 0xD2, 0xFE),
          align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print("wrote", OUT, "—", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    build()
